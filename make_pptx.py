import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_cetam_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Cores CETAM
    BG_COLOR = RGBColor(6, 14, 26)
    CARD_BG = RGBColor(12, 26, 48)
    CARD_BORDER = RGBColor(30, 58, 95)
    GREEN = RGBColor(0, 208, 98)
    GOLD = RGBColor(255, 184, 28)
    BLUE = RGBColor(56, 189, 248)
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(203, 213, 225)
    RED = RGBColor(239, 68, 68)

    def set_slide_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()
        return bg

    def add_top_header(slide, title_text, presenter_text=None):
        # Header Box
        header = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.65))
        header.fill.solid()
        header.fill.fore_color.rgb = CARD_BG
        header.line.color.rgb = CARD_BORDER
        header.line.width = Pt(1)

        tf = header.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.1)

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = WHITE

        if presenter_text:
            # Presenter Badge
            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(0.43), Inches(2.8), Inches(0.48))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(25, 40, 65)
            badge.line.color.rgb = GOLD
            badge.line.width = Pt(1)
            btf = badge.text_frame
            btf.margin_top = Inches(0.06)
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            r1 = bp.add_run()
            r1.text = "Apresentador: "
            r1.font.size = Pt(12)
            r1.font.color.rgb = GOLD
            r2 = bp.add_run()
            r2.text = presenter_text
            r2.font.size = Pt(13)
            r2.font.bold = True
            r2.font.color.rgb = WHITE

    def add_column_cards(slide, col_index, step1_desc, step2_desc, step3_cmd, step3_res, step3_is_error=False):
        col_width = Inches(3.95)
        col_gap = Inches(0.24)
        left = Inches(0.5) + col_index * (col_width + col_gap)

        # Card 1: Precisa ser feito
        c1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.2), col_width, Inches(1.6))
        c1.fill.solid()
        c1.fill.fore_color.rgb = CARD_BG
        c1.line.color.rgb = BLUE
        c1.line.width = Pt(2)
        tf1 = c1.text_frame
        tf1.word_wrap = True
        tf1.margin_left = Inches(0.15)
        tf1.margin_top = Inches(0.1)
        p1 = tf1.paragraphs[0]
        p1.text = "👤 1. Precisa ser feito"
        p1.font.size = Pt(13)
        p1.font.bold = True
        p1.font.color.rgb = BLUE
        p1_sub = tf1.add_paragraph()
        p1_sub.text = step1_desc
        p1_sub.font.size = Pt(12)
        p1_sub.font.color.rgb = WHITE

        # Arrow 1
        a1 = slide.shapes.add_textbox(left, Inches(2.82), col_width, Inches(0.3))
        a1.text_frame.paragraphs[0].text = "↓"
        a1.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        a1.text_frame.paragraphs[0].font.size = Pt(14)
        a1.text_frame.paragraphs[0].font.color.rgb = MUTED

        # Card 2: Como fazer
        c2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.15), col_width, Inches(2.2))
        c2.fill.solid()
        c2.fill.fore_color.rgb = CARD_BG
        c2.line.color.rgb = GREEN
        c2.line.width = Pt(2)
        tf2 = c2.text_frame
        tf2.word_wrap = True
        tf2.margin_left = Inches(0.15)
        tf2.margin_top = Inches(0.1)
        p2 = tf2.paragraphs[0]
        p2.text = "⚙️ 2. Como fazer"
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = GREEN
        p2_sub = tf2.add_paragraph()
        p2_sub.text = step2_desc
        p2_sub.font.size = Pt(11.5)
        p2_sub.font.color.rgb = WHITE

        # Arrow 2
        a2 = slide.shapes.add_textbox(left, Inches(5.37), col_width, Inches(0.3))
        a2.text_frame.paragraphs[0].text = "↓"
        a2.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        a2.text_frame.paragraphs[0].font.size = Pt(14)
        a2.text_frame.paragraphs[0].font.color.rgb = MUTED

        # Card 3: Como testar
        c3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(5.68), col_width, Inches(1.4))
        c3.fill.solid()
        c3.fill.fore_color.rgb = CARD_BG
        c3.line.color.rgb = GOLD
        c3.line.width = Pt(2)
        tf3 = c3.text_frame
        tf3.word_wrap = True
        tf3.margin_left = Inches(0.15)
        tf3.margin_top = Inches(0.08)
        p3 = tf3.paragraphs[0]
        p3.text = "✉️ 3. Como testar"
        p3.font.size = Pt(13)
        p3.font.bold = True
        p3.font.color.rgb = GOLD
        if step3_cmd:
            p3_cmd = tf3.add_paragraph()
            p3_cmd.text = f"CMD/PowerShell: {step3_cmd}"
            p3_cmd.font.size = Pt(11)
            p3_cmd.font.bold = True
            p3_cmd.font.color.rgb = BLUE
        p3_res = tf3.add_paragraph()
        p3_res.text = step3_res
        p3_res.font.size = Pt(11)
        p3_res.font.color.rgb = RED if step3_is_error else GREEN

    # ==========================================
    # SLIDE 1: CAPA
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1)

    # Header Institucional
    add_top_header(s1, "🏛️ CETAM • Centro de Educação Tecnológica do Amazonas", "Equipe 3")

    # Title Hero
    tb = s1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(2.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run()
    r1.text = "Serviços Corporativos e Segurança no "
    r1.font.size = Pt(32)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = p.add_run()
    r2.text = "Windows Server 2022\n"
    r2.font.size = Pt(32)
    r2.font.bold = True
    r2.font.color.rgb = GREEN

    p_sub = tf.add_paragraph()
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "Implementação de Diretivas de Grupo (GPO), Servidor de Arquivos, Fila de Impressão e Firewall"
    p_sub.font.size = Pt(16)
    p_sub.font.color.rgb = MUTED

    # 5 Integrantes Cards
    members = [
        ("Elizandra", "Abertura & Apresentação", "Introdução do projeto, alinhamento e coordenação."),
        ("Gilvan", "Diretivas de Grupo (GPO)", "Senhas fortes, bloqueio USB, CMD e exceção TI."),
        ("Mariane", "Servidor de Arquivos", "8 pastas departamentais, NTFS e Acesso Negado."),
        ("Edmar", "Servidor de Impressão", "Fila IMP_Geral (192.168.10.20) e correção RPC."),
        ("Adler", "Firewall & DNS", "Perfis ativos, liberação SMB/Ping e DNS Reverso.")
    ]

    card_w = Inches(2.3)
    card_gap = Inches(0.2)
    for i, (name, role, desc) in enumerate(members):
        left = Inches(0.5) + i * (card_w + card_gap)
        mc = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(4.2), card_w, Inches(2.6))
        mc.fill.solid()
        mc.fill.fore_color.rgb = CARD_BG
        mc.line.color.rgb = CARD_BORDER
        mc.line.width = Pt(1)
        mtf = mc.text_frame
        mtf.word_wrap = True
        mtf.margin_left = Inches(0.1)
        mtf.margin_top = Inches(0.2)

        mp1 = mtf.paragraphs[0]
        mp1.alignment = PP_ALIGN.CENTER
        mp1.text = name
        mp1.font.size = Pt(18)
        mp1.font.bold = True
        mp1.font.color.rgb = WHITE

        mp2 = mtf.add_paragraph()
        mp2.alignment = PP_ALIGN.CENTER
        mp2.text = role
        mp2.font.size = Pt(12)
        mp2.font.bold = True
        mp2.font.color.rgb = GOLD

        mp3 = mtf.add_paragraph()
        mp3.alignment = PP_ALIGN.CENTER
        mp3.text = desc
        mp3.font.size = Pt(11)
        mp3.font.color.rgb = MUTED

    # ==========================================
    # SLIDE 2: GILVAN - GPO
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2)
    add_top_header(s2, "🛡️ Gestão de Identidade e Políticas de Grupo (GPO)", "Gilvan")

    add_column_cards(s2, 0,
        "Exigir senhas fortes de 8 dígitos e bloquear a conta após 5 tentativas incorretas.",
        "gpmc.msc ➔ Domain ➔ Default Domain Policy ➔ Edit ➔ Computer Configuration ➔ Policies ➔ Windows Settings ➔ Security Settings ➔ Account Policies (Password & Account Lockout Policy).",
        "net accounts",
        "✓ Política auditada e ativa no domínio.")

    add_column_cards(s2, 1,
        "Bloquear pendrives USB e Prompt de Comando (CMD/PowerShell) para evitar vazamento ou execução indevida.",
        "gpmc.msc ➔ GPO_Seguranca_e_Padronizacao ➔ Edit:\n• USB: Computer Config ➔ Admin Templates ➔ System ➔ Removable Storage Access ➔ Deny all access.\n• CMD: User Config ➔ Admin Templates ➔ System ➔ Prevent access to the command prompt ➔ Enabled.",
        None,
        "🚫 Pop-up: 'Access is denied' / 'The command prompt has been disabled by your administrator.'",
        step3_is_error=True)

    add_column_cards(s2, 2,
        "Fixar wallpaper institucional e garantir acesso total para a equipe de TI (Domain Admins).",
        "gpmc.msc ➔ Group Policy Objects:\n• Wallpaper: Right-click em GPO_Seguranca_e_Padronizacao ➔ Edit ➔ User Config ➔ Policies ➔ Admin Templates ➔ Desktop ➔ Desktop ➔ Desktop Wallpaper: \\\\SRV-DC01\\Dados\\Publico\\wallpaper.jpg\n• Exceção TI: Clicar em GPO_Seguranca_e_Padronizacao ➔ Aba Delegation ➔ Advanced ➔ Domain Admins ➔ Desmarcar 'Apply group policy'.",
        "gpresult /r",
        "✓ GPO aplicada a usuários e ignorada para TI.")

    # ==========================================
    # SLIDE 3: MARIANE - ARQUIVOS
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3)
    add_top_header(s3, "📁 Servidor de Arquivos e Permissões Departamentais", "Mariane")

    add_column_cards(s3, 0,
        "Criar as 8 pastas departamentais em C:\\Dados e compartilhá-las como \\\\esperancadigital.local\\Dados.",
        "• Pastas: Criar em C:\\Dados (Diretoria, Coordenacao, Secretaria, Biblioteca, Recepcao, Professores, Alunos, Publico).\n• Share: Server Manager ➔ File and Storage Services ➔ Shares ➔ New Share (SMB - Quick) ➔ Path: C:\\Dados ➔ Share Name: 'Dados'.",
        "\\\\esperancadigital.local\\Dados",
        "✓ Todas as 8 pastas visíveis na rede.")

    add_column_cards(s3, 1,
        "Permitir que cada usuário leia, crie e edite arquivos apenas na pasta do seu departamento.",
        "Right-click na pasta (ex: Secretaria) ➔ Properties ➔ Security ➔ Advanced ➔ Disable inheritance ('Convert...') ➔ Remove Users ➔ Add: 'GG_Secretaria' ➔ Permission: Modify.",
        "\\\\esperancadigital.local\\Dados\\Secretaria",
        "✓ Acesso liberado para ler, criar e salvar.")

    add_column_cards(s3, 2,
        "Bloquear o acesso a pastas de outros setores exibindo o pop-up de Acesso Negado.",
        "Server Manager ➔ File and Storage Services ➔ Shares ➔ Right-click 'Dados' ➔ Properties ➔ Settings ➔ Uncheck 'Enable access-based enumeration'.",
        "\\\\esperancadigital.local\\Dados\\Diretoria",
        "🚫 Pop-up na tela: 'Network Error: Access is denied.'",
        step3_is_error=True)

    # ==========================================
    # SLIDE 4: EDMAR - IMPRESSÃO
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4)
    add_top_header(s4, "🖨️ Servidor de Impressão e Fila Geral Centralizada", "Edmar")

    add_column_cards(s4, 0,
        "Habilitar o serviço central de impressão no Windows Server 2022.",
        "Server Manager ➔ Manage ➔ Add Roles and Features ➔ Server Roles ➔ Print and Document Services ➔ Role Services: Print Server ➔ Install.",
        "printmanagement.msc",
        "✓ Print Servers ➔ SRV-DC01 ativo e operacional.")

    add_column_cards(s4, 1,
        "Criar a fila corporativa IMP_Geral no IP 192.168.10.20 publicada no AD.",
        "printmanagement.msc ➔ Print Servers ➔ SRV-DC01 ➔ Right-click 'Printers' ➔ Add Printer... ➔ TCP/IP (192.168.10.20) ➔ Nome: IMP_Geral ➔ Marcar 'Share this printer' e 'List in the directory'.",
        "\\\\SRV-DC01\\IMP_Geral",
        "✓ Fila instalada e pronta para todos os setores.")

    add_column_cards(s4, 2,
        "Corrigir erro de conexão remota 0x8007011b para permitir impressão direta das estações.",
        "PowerShell (Admin): New-ItemProperty -Path 'HKLM:\\System\\CurrentControlSet\\Control\\Print' -Name 'RpcAuthnLevelPrivacyEnabled' -Value 0 -PropertyType DWORD -Force; Restart-Service Spooler.",
        "IMP_Geral ➔ Print Test Page",
        "✓ Página de teste impressa sem travar.")

    # ==========================================
    # SLIDE 5: ADLER - FIREWALL & DNS
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5)
    add_top_header(s5, "🛡️ Windows Defender Firewall & Conectividade de Rede", "Adler")

    add_column_cards(s5, 0,
        "Manter o Firewall ativo nos 3 perfis (Domain, Private e Public) para segurança da rede.",
        "wf.msc ➔ Windows Defender Firewall Properties ➔ Abas 'Domain', 'Private' e 'Public Profile' ➔ Firewall state: 'On (recommended)'.",
        "wf.msc",
        "✓ Status verde ativo em Domain, Private e Public.")

    add_column_cards(s5, 1,
        "Liberar o tráfego para compartilhamento SMB (Porta 445) e resposta a testes de Ping.",
        "wf.msc ➔ Inbound Rules ➔ Habilitar (Enable Rule):\n• 'File and Printer Sharing (SMB-In)'\n• 'File and Printer Sharing (Echo Request - ICMPv4-In)'.",
        "ping 192.168.10.2",
        "✓ 0% de perda com latência <1ms.")

    add_column_cards(s5, 2,
        "Criar a resolução de DNS Reverso para resolver o IP 192.168.10.2 e evitar lentidão no logon.",
        "dnsmgmt.msc ➔ SRV-DC01 ➔ Reverse Lookup Zones ➔ New Zone... ➔ Primary (IPv4: 192.168.10) ➔ New Pointer (PTR): IP .2 ➔ SRV-DC01.esperancadigital.local.",
        "nslookup 192.168.10.2",
        "✓ Retorna: SRV-DC01.esperancadigital.local.")

    # ==========================================
    # SLIDE 6: ENCERRAMENTO
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6)

    # Topo Institucional
    add_top_header(s6, "🏛️ CETAM • Centro de Educação Tecnológica do Amazonas", "Elizandra & Equipe")

    # Centro
    tb6 = s6.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.0))
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    p6_icon = tf6.paragraphs[0]
    p6_icon.alignment = PP_ALIGN.CENTER
    p6_icon.text = "🏛️"
    p6_icon.font.size = Pt(36)

    p6_title = tf6.add_paragraph()
    p6_title.alignment = PP_ALIGN.CENTER
    r6_1 = p6_title.add_run()
    r6_1.text = "Muito Obrigado pela "
    r6_1.font.size = Pt(38)
    r6_1.font.bold = True
    r6_1.font.color.rgb = WHITE
    r6_2 = p6_title.add_run()
    r6_2.text = "Atenção!\n"
    r6_2.font.size = Pt(38)
    r6_2.font.bold = True
    r6_2.font.color.rgb = GREEN

    p6_sub = tf6.add_paragraph()
    p6_sub.alignment = PP_ALIGN.CENTER
    p6_sub.text = "Administração, Segurança e Serviços Corporativos no Windows Server 2022"
    p6_sub.font.size = Pt(18)
    p6_sub.font.color.rgb = MUTED

    # Bottom References
    tb_ref = s6.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.333), Inches(0.8))
    tf_ref = tb_ref.text_frame
    tf_ref.word_wrap = True
    p_ref = tf_ref.paragraphs[0]
    p_ref.alignment = PP_ALIGN.CENTER
    p_ref.text = "📚 Referências & Suporte: Microsoft Learn (GPMC) • Microsoft TechNet (SMB/NTFS) • Microsoft KB5005652 (RPC) • Roteiro Técnico CETAM • Google Gemini 3.7 & Antigravity IDE"
    p_ref.font.size = Pt(12)
    p_ref.font.color.rgb = MUTED

    # Salvar arquivo
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Apresentacao_Grupo3_CETAM.pptx")
    prs.save(output_path)
    print(f"[+] Sucesso! Arquivo PowerPoint gerado em: {output_path}")

if __name__ == "__main__":
    create_cetam_presentation()
